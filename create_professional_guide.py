#!/usr/bin/env python3
"""
Script to create professional documentation from the RASS Academy User Guide
Generates DOCX and PPTX versions with proper formatting
"""

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import re

def create_docx_guide():
    """Create a professional DOCX document from the markdown guide"""
    print("Creating DOCX document...")
    
    # Create document
    doc = Document()
    
    # Set document margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
    
    # Read the markdown content
    with open('RASS_ACADEMY_COMPLETE_USER_GUIDE.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Add title page
    title = doc.add_heading('RASS ACADEMY', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.runs[0]
    title_run.font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle = doc.add_heading('Complete Learning Management System', level=2)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.runs[0]
    subtitle_run.font.color.rgb = RGBColor(0, 102, 204)
    
    subtitle2 = doc.add_heading('User Guide & Documentation', level=2)
    subtitle2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()
    doc.add_paragraph()
    
    # Add version info
    version_para = doc.add_paragraph()
    version_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    version_para.add_run('Version 1.0 | 2025\n').bold = True
    version_para.add_run('Comprehensive Guide for Students, Instructors, and Administrators')
    
    doc.add_page_break()
    
    # Process markdown content
    lines = content.split('\n')
    skip_yaml = False
    
    for line in lines:
        # Skip YAML front matter
        if line.strip() == '---':
            skip_yaml = not skip_yaml
            continue
        if skip_yaml:
            continue
            
        # Skip the first title (already added)
        if line.startswith('# 📚 RASS ACADEMY'):
            continue
        if line.startswith('## Complete Learning Management System User Guide'):
            continue
            
        # Process headers
        if line.startswith('####'):
            heading = line.replace('####', '').strip()
            doc.add_heading(heading, level=4)
        elif line.startswith('###'):
            heading = line.replace('###', '').strip()
            doc.add_heading(heading, level=3)
        elif line.startswith('##'):
            heading = line.replace('##', '').strip()
            doc.add_heading(heading, level=2)
        elif line.startswith('#'):
            heading = line.replace('#', '').strip()
            if heading:  # Skip empty headers
                doc.add_heading(heading, level=1)
        
        # Process bullet points
        elif line.strip().startswith('- '):
            text = line.strip()[2:]
            para = doc.add_paragraph(text, style='List Bullet')
        elif line.strip().startswith('* '):
            text = line.strip()[2:]
            para = doc.add_paragraph(text, style='List Bullet')
        
        # Process numbered lists
        elif re.match(r'^\d+\.', line.strip()):
            text = re.sub(r'^\d+\.\s*', '', line.strip())
            para = doc.add_paragraph(text, style='List Number')
        
        # Process code blocks
        elif line.strip().startswith('```'):
            continue  # Skip code fence markers
        
        # Process horizontal rules
        elif line.strip() == '---':
            doc.add_paragraph('_' * 50)
        
        # Process regular text
        elif line.strip() and not line.startswith('**'):
            # Handle bold text
            if '**' in line:
                para = doc.add_paragraph()
                parts = line.split('**')
                for i, part in enumerate(parts):
                    if i % 2 == 0:
                        para.add_run(part)
                    else:
                        para.add_run(part).bold = True
            else:
                doc.add_paragraph(line.strip())
    
    # Save document
    doc.save('RASS_ACADEMY_USER_GUIDE.docx')
    print("✓ DOCX document created: RASS_ACADEMY_USER_GUIDE.docx")

def create_pptx_presentation():
    """Create a professional PowerPoint presentation"""
    print("Creating PowerPoint presentation...")
    
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    # Define color scheme
    DARK_BLUE = PptxRGBColor(0, 51, 102)
    LIGHT_BLUE = PptxRGBColor(0, 102, 204)
    ACCENT = PptxRGBColor(255, 140, 0)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = PptxInches(1)
    top = PptxInches(2.5)
    width = PptxInches(8)
    height = PptxInches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "RASS ACADEMY"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = PptxPt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(left, top + PptxInches(1), width, PptxInches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Complete Learning Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = PptxPt(32)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    
    # Add description
    desc_box = slide.shapes.add_textbox(left, top + PptxInches(1.8), width, PptxInches(0.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Comprehensive User Guide & Visual Workflows"
    desc_para = desc_frame.paragraphs[0]
    desc_para.alignment = PP_ALIGN.CENTER
    desc_para.font.size = PptxPt(24)
    desc_para.font.color.rgb = ACCENT
    
    # Slide 2: Table of Contents
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📖 Table of Contents"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Introduction to RASS Academy"
    
    for item in [
        "2. Getting Started",
        "3. User Roles & Dashboards",
        "4. Student Portal Features",
        "5. Instructor Portal Features",
        "6. Administrator Portal Features",
        "7. Public Features & Partnerships",
        "8. Technical Setup & Deployment",
        "9. Troubleshooting & Support"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(20)
        p.level = 0
    
    # Slide 3: What is RASS Academy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎓 What is RASS Academy?"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A comprehensive Learning Management System (LMS) designed to deliver modern, interactive online education"
    tf.paragraphs[0].font.size = PptxPt(20)
    
    for item in [
        "🎓 Comprehensive Course Management",
        "👥 Batch-based Learning",
        "📊 Advanced Analytics & Progress Tracking",
        "🎥 Video-based Learning with Interactive Features",
        "📝 Assessments & Auto-grading",
        "🏆 Gamification & Leaderboards",
        "💬 Real-time Communication Tools",
        "🤝 Partnership Programs"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18)
        p.level = 0
    
    # Slide 4: Key Features for Students
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "👨‍🎓 Student Features"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Access to structured courses with video lectures"
    
    for item in [
        "Interactive quizzes with instant feedback",
        "Assignment submission and tracking",
        "Live sessions and discussions",
        "AI-powered doubt clearing",
        "Progress tracking and certificates",
        "Leaderboards and achievements",
        "Timestamped video notes",
        "Discussion forums"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18)
        p.level = 0
    
    # Slide 5: Key Features for Instructors
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "👨‍🏫 Instructor Features"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Course and content management"
    
    for item in [
        "Batch management and analytics",
        "Student performance tracking",
        "Assignment and quiz creation",
        "Attendance management (manual + auto)",
        "Announcements and communication",
        "Advanced analytics dashboard",
        "Live session hosting",
        "Grading and feedback tools"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18)
        p.level = 0
    
    # Slide 6: Key Features for Admins
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "⚙️ Administrator Features"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Complete platform management"
    
    for item in [
        "User management (students, instructors)",
        "Course and batch administration",
        "Certificate generation and management",
        "Partnership program oversight",
        "Platform analytics and reporting",
        "Event management",
        "Support ticket system",
        "System configuration"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18)
        p.level = 0
    
    # Slide 7: Student Learning Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎯 Student Learning Workflow"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Browse Course Catalog → Filter & Search"
    
    for item in [
        "2. View Course Details → Check Curriculum & Reviews",
        "3. Enroll in Course → Complete Payment",
        "4. Access Course → Watch Video Lectures",
        "5. Take Notes → Timestamp-based annotations",
        "6. Complete Assignments → Submit & Track",
        "7. Take Quizzes → Auto-graded assessments",
        "8. Join Live Sessions → Real-time interaction",
        "9. Ask Doubts → AI + Instructor help",
        "10. Earn Certificate → 100% completion + pass all tests"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    # Slide 8: Instructor Course Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📚 Instructor Course Management"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Create Course → Add Basic Information"
    
    for item in [
        "2. Build Curriculum → Add Modules & Lessons",
        "3. Upload Content → Videos, PDFs, Resources",
        "4. Create Assessments → Assignments & Quizzes",
        "5. Set up Batch → Assign Students & Schedule",
        "6. Publish Course → Students Can Enroll",
        "7. Monitor Progress → Analytics Dashboard",
        "8. Grade Submissions → Provide Feedback",
        "9. Respond to Doubts → Help Students",
        "10. Generate Reports → Track Performance"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    # Slide 9: Assessment System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📝 Assessment System"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Assignments"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = PptxPt(20)
    
    for item in [
        "• File upload (PDF, DOC, ZIP)",
        "• Manual grading with rubrics",
        "• Feedback and comments",
        "• Late submission tracking"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Quizzes"
    p.font.bold = True
    p.font.size = PptxPt(20)
    
    for item in [
        "• Multiple choice, True/False, Coding problems",
        "• Auto-grading with instant feedback",
        "• Timed assessments",
        "• Multiple attempts tracking"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    # Slide 10: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "💻 Technology Stack"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = PptxPt(20)
    
    for item in ["React 18 + TypeScript", "Vite Build Tool", "TailwindCSS + Radix UI"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Backend"
    p.font.bold = True
    p.font.size = PptxPt(20)
    
    for item in ["Node.js + Express", "MongoDB + Mongoose", "JWT Authentication"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Integrations"
    p.font.bold = True
    p.font.size = PptxPt(20)
    
    for item in ["Razorpay (Payments)", "Nodemailer (Email)", "reCAPTCHA (Security)"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(16)
        p.level = 0
    
    # Slide 11: Support & Help
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🆘 Support & Help"
    title.text_frame.paragraphs[0].font.size = PptxPt(40)
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Email: support@rassacademy.com"
    tf.paragraphs[0].font.size = PptxPt(20)
    
    for item in [
        "Phone: +XX-XXXX-XXXX (Mon-Fri, 10 AM - 5 PM)",
        "Live Chat: Available on website (business hours)",
        "Support Tickets: Create from dashboard",
        "Help Center: Searchable knowledge base",
        "Community Forum: Peer support",
        "Video Tutorials: Step-by-step guides",
        "Social Media: @rassacademy"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18)
        p.level = 0
    
    # Final Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add closing message
    left = PptxInches(2)
    top = PptxInches(3)
    width = PptxInches(6)
    height = PptxInches(1.5)
    
    closing_box = slide.shapes.add_textbox(left, top, width, height)
    closing_frame = closing_box.text_frame
    closing_frame.text = "Thank You!"
    closing_para = closing_frame.paragraphs[0]
    closing_para.alignment = PP_ALIGN.CENTER
    closing_para.font.size = PptxPt(48)
    closing_para.font.bold = True
    closing_para.font.color.rgb = DARK_BLUE
    
    # Add contact info
    contact_box = slide.shapes.add_textbox(left, top + PptxInches(1.2), width, PptxInches(0.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "www.rassacademy.com | support@rassacademy.com"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = PptxPt(20)
    contact_para.font.color.rgb = LIGHT_BLUE
    
    # Save presentation
    prs.save('RASS_ACADEMY_USER_GUIDE.pptx')
    print("✓ PowerPoint presentation created: RASS_ACADEMY_USER_GUIDE.pptx")

if __name__ == "__main__":
    print("=" * 60)
    print("RASS Academy Professional Documentation Generator")
    print("=" * 60)
    print()
    
    create_docx_guide()
    print()
    create_pptx_presentation()
    print()
    print("=" * 60)
    print("✓ All documents created successfully!")
    print("=" * 60)
    print()
    print("Generated files:")
    print("  1. RASS_ACADEMY_COMPLETE_USER_GUIDE.md (Markdown)")
    print("  2. RASS_ACADEMY_USER_GUIDE.docx (Word Document)")
    print("  3. RASS_ACADEMY_USER_GUIDE.pptx (PowerPoint)")
    print()
